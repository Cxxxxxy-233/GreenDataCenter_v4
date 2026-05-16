from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
ASSET_DIR = OUT_DIR / "assets"
PPTX_PATH = OUT_DIR / "AI_Green_Data_Center_Planning_editable_v2.pptx"

SLIDE_W = 13.333333
SLIDE_H = 7.5
PX_W = 2752
PX_H = 1536

FONT = "Microsoft YaHei"

GREEN = RGBColor(30, 188, 149)
GREEN_DARK = RGBColor(10, 166, 126)
GREEN_LIGHT = RGBColor(233, 251, 245)
GREEN_LINE = RGBColor(76, 217, 181)
GREEN_ACCENT = RGBColor(110, 229, 193)
TEXT = RGBColor(24, 31, 31)
TEXT_SOFT = RGBColor(70, 78, 78)
GRID = RGBColor(234, 242, 240)
WHITE = RGBColor(255, 255, 255)
CYAN = RGBColor(77, 214, 233)


def px_to_in(x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    return x / PX_W * SLIDE_W, y / PX_H * SLIDE_H, w / PX_W * SLIDE_W, h / PX_H * SLIDE_H


def add_text(
    slide,
    text: str,
    box_px: tuple[float, float, float, float],
    font_size: float,
    *,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    wrap: bool = True,
) -> None:
    x, y, w, h = px_to_in(*box_px)
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.vertical_anchor = vertical
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_round_rect(
    slide,
    box_px: tuple[float, float, float, float],
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = GREEN_LINE,
    line_width: float = 1.3,
    shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    transparency: float = 0.0,
) -> None:
    x, y, w, h = px_to_in(*box_px)
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)


def add_circle(slide, box_px: tuple[float, float, float, float], *, fill: RGBColor, line: RGBColor = GREEN_LINE) -> None:
    add_round_rect(slide, box_px, fill=fill, line=line, shape_type=MSO_AUTO_SHAPE_TYPE.OVAL)


def add_line(slide, start_px: tuple[float, float], end_px: tuple[float, float], *, color: RGBColor = GREEN_LINE, width: float = 1.0) -> None:
    x1 = start_px[0] / PX_W * SLIDE_W
    y1 = start_px[1] / PX_H * SLIDE_H
    x2 = end_px[0] / PX_W * SLIDE_W
    y2 = end_px[1] / PX_H * SLIDE_H
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)


def add_picture_contain(slide, path: Path, box_px: tuple[float, float, float, float]) -> None:
    with Image.open(path) as img:
        iw, ih = img.size
    bx, by, bw, bh = box_px
    src_ratio = iw / ih
    box_ratio = bw / bh
    if src_ratio > box_ratio:
        fw = bw
        fh = fw / src_ratio
        fx = bx
        fy = by + (bh - fh) / 2
    else:
        fh = bh
        fw = fh * src_ratio
        fy = by
        fx = bx + (bw - fw) / 2
    x, y, w, h = px_to_in(fx, fy, fw, fh)
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def alpha_crop(page_name: str, crop_box: tuple[int, int, int, int], out_name: str, threshold: int = 245) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    img = Image.open(OUT_DIR / page_name).convert("RGBA")
    cropped = img.crop(crop_box)
    pixels = list(cropped.getdata())
    result = []
    for r, g, b, a in pixels:
        avg = (r + g + b) // 3
        if r > threshold and g > threshold and b > threshold:
            result.append((255, 255, 255, 0))
        elif avg > 228:
            alpha = max(0, min(255, int((255 - avg) * 6)))
            result.append((r, g, b, alpha))
        else:
            result.append((r, g, b, a))
    cropped.putdata(result)
    out = ASSET_DIR / out_name
    cropped.save(out)
    return out


def build_assets() -> dict[str, Path]:
    return {
        "p1_person_1": alpha_crop("page_01.png", (610, 360, 930, 910), "p1_person_1.png"),
        "p1_person_2": alpha_crop("page_01.png", (1420, 390, 1810, 900), "p1_person_2.png"),
        "p1_person_3": alpha_crop("page_01.png", (2330, 380, 2750, 920), "p1_person_3.png"),
        "p1_gauge_1": alpha_crop("page_01.png", (255, 1040, 505, 1302), "p1_gauge_1.png"),
        "p1_gauge_2": alpha_crop("page_01.png", (1128, 1040, 1376, 1302), "p1_gauge_2.png"),
        "p1_gauge_3": alpha_crop("page_01.png", (1998, 1040, 2248, 1302), "p1_gauge_3.png"),
        "p3_icon_1": alpha_crop("page_03.png", (92, 350, 355, 618), "p3_icon_1.png"),
        "p3_icon_2": alpha_crop("page_03.png", (1002, 350, 1265, 618), "p3_icon_2.png"),
        "p3_icon_3": alpha_crop("page_03.png", (1868, 350, 2148, 618), "p3_icon_3.png"),
        "p3_cycle": alpha_crop("page_03.png", (1580, 930, 2130, 1360), "p3_cycle.png"),
        "p3_gavel": alpha_crop("page_03.png", (2390, 900, 2665, 1185), "p3_gavel.png"),
    }


def add_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    for x in range(0, PX_W, 210):
        add_line(slide, (x, 0), (x, PX_H), color=GRID, width=0.6)
    for y in range(0, PX_H, 210):
        add_line(slide, (0, y), (PX_W, y), color=GRID, width=0.6)


def slide1(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_text(
        slide,
        "数据中心绿电一体化方案智能规划系统",
        (570, 50, 1800, 120),
        28,
        color=GREEN_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
        wrap=False,
    )
    add_text(
        slide,
        "基于 LangChain + LangGraph 的多智能体协同架构",
        (700, 175, 1420, 72),
        18,
        align=PP_ALIGN.CENTER,
        wrap=False,
    )

    cards = [
        {
            "box": (120, 345, 805, 530),
            "title": "能耗指数级增长",
            "number": "2166亿",
            "unit": "千瓦时",
            "desc": "我国数据中心用电量持续攀升",
            "asset": assets["p1_person_1"],
            "asset_box": (630, 430, 230, 315),
        },
        {
            "box": (972, 345, 805, 530),
            "title": "碳排压力显著",
            "number": "1350万",
            "unit": "吨",
            "desc": "年碳排放量巨大，刚性指标收紧",
            "asset": assets["p1_person_2"],
            "asset_box": (1495, 420, 250, 320),
        },
        {
            "box": (1825, 345, 805, 530),
            "title": "传统模式滞后",
            "number": "45%",
            "unit": "",
            "desc": "高度依赖人工经验，设计流程复杂且协同成本高",
            "asset": assets["p1_person_3"],
            "asset_box": (2330, 420, 250, 320),
        },
    ]
    for item in cards:
        x, y, w, h = item["box"]
        add_round_rect(slide, item["box"], fill=WHITE, line=GREEN_LINE, line_width=1.8)
        add_round_rect(slide, (x + 22, y + 20, w - 44, h - 44), fill=GREEN_LIGHT, line=GREEN_LIGHT, line_width=0.5, transparency=0.08)
        add_text(slide, item["title"], (x + 55, y + 55, 330, 84), 20, bold=True)
        add_text(slide, item["number"], (x + 60, y + 160, 330, 92), 43, color=GREEN_DARK, bold=True, wrap=False)
        if item["unit"]:
            add_text(slide, item["unit"], (x + 250, y + 252, 120, 32), 15, bold=True, color=TEXT_SOFT, wrap=False)
        add_text(slide, item["desc"], (x + 60, y + 380, 430, 82), 17, color=TEXT_SOFT)
        add_picture_contain(slide, item["asset"], item["asset_box"])

    add_round_rect(slide, (950, 935, 855, 120), fill=GREEN_DARK, line=GREEN_DARK, line_width=1.0)
    add_text(
        slide,
        "算力基建高耗能危机与双碳战略需求",
        (985, 960, 785, 64),
        20,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    gauges = [
        {
            "asset": assets["p1_gauge_1"],
            "asset_box": (255, 1088, 205, 205),
            "value": "±1%",
            "title": "PUE（电能利用效率）",
            "desc": "极值挑战",
            "text_box": (500, 1095, 310, 96),
        },
        {
            "asset": assets["p1_gauge_2"],
            "asset_box": (1128, 1088, 205, 205),
            "value": "30%",
            "title": "CUE（碳效指标）",
            "desc": "刚性约束",
            "text_box": (1370, 1095, 285, 96),
        },
        {
            "asset": assets["p1_gauge_3"],
            "asset_box": (2005, 1088, 205, 205),
            "value": "95%",
            "title": "全生命周期 ROI 与经济性",
            "desc": "",
            "text_box": (2245, 1088, 260, 110),
        },
    ]
    for item in gauges:
        add_picture_contain(slide, item["asset"], item["asset_box"])
        ax, ay, aw, ah = item["asset_box"]
        add_text(slide, item["value"], (ax + 44, ay + 134, 118, 36), 18, bold=True, align=PP_ALIGN.CENTER)
        tx, ty, tw, th = item["text_box"]
        add_text(slide, item["title"], (tx, ty, tw, 44), 15, bold=False)
        if item["desc"]:
            add_text(slide, item["desc"], (tx, ty + 38, tw, 34), 14, color=TEXT_SOFT)

    add_round_rect(slide, (145, 1385, 2460, 104), fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.0)
    add_text(
        slide,
        "以 AI 多智能体协同，突破传统人工规划瓶颈，实现数据中心建设 100% 数据驱动的低碳论证。",
        (240, 1410, 2270, 50),
        16,
        align=PP_ALIGN.CENTER,
        color=TEXT_SOFT,
    )


def slide2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_text(slide, "系统整体技术架构蓝图", (160, 55, 900, 88), 27, bold=True)
    add_text(slide, "5 层核心架构与全链路数据流向", (160, 150, 820, 56), 18, color=TEXT_SOFT)

    add_line(slide, (220, 275), (220, 1435), color=GREEN_DARK, width=2.8)
    add_text(slide, "用户参数", (110, 760, 120, 220), 18, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "↓", (190, 695, 60, 60), 34, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "↓", (190, 1290, 60, 60), 34, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

    add_line(slide, (2585, 275), (2585, 1435), color=CYAN, width=2.8)
    add_text(slide, "最优推荐方案", (2488, 760, 140, 240), 18, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "↑", (2555, 300, 60, 60), 34, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    # Layer 1
    add_round_rect(slide, (330, 275, 2130, 180), fill=WHITE, line=GRID, line_width=1.0)
    add_text(slide, "第1层：前端展示层", (395, 307, 290, 38), 18, bold=True)
    add_round_rect(slide, (1100, 285, 560, 76), fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.8)
    add_text(slide, "Vue 3 + Vite + Element Plus", (1125, 302, 510, 40), 20, bold=True, align=PP_ALIGN.CENTER, wrap=False)
    front_items = [
        (395, 382, 420, 56, "参数配置模块"),
        (915, 382, 420, 56, "方案生成模块"),
        (1435, 382, 420, 56, "结果展示模块"),
        (1955, 382, 420, 56, "历史记录模块"),
    ]
    for x, y, w, h, label in front_items:
        add_round_rect(slide, (x, y, w, h), fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.2)
        add_text(slide, label, (x + 20, y + 8, w - 40, 38), 18, bold=True, align=PP_ALIGN.CENTER)

    # Layer 2
    add_round_rect(slide, (330, 495, 2130, 165), fill=WHITE, line=GRID, line_width=1.0)
    add_text(slide, "第2层：后端服务层", (395, 527, 290, 38), 18, bold=True)
    add_round_rect(slide, (1100, 500, 560, 76), fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.8)
    add_text(slide, "FastAPI + Python | PostgreSQL", (1120, 515, 520, 42), 20, bold=True, align=PP_ALIGN.CENTER, wrap=False)
    for x, y, w, h, label in [
        (515, 590, 450, 52, "RESTful API"),
        (1130, 590, 420, 52, "参数校验"),
        (1710, 590, 500, 52, "数据持久化"),
    ]:
        add_round_rect(slide, (x, y, w, h), fill=WHITE, line=GRID, line_width=1.0)
        add_text(slide, label, (x + 20, y + 7, w - 40, 36), 17, align=PP_ALIGN.CENTER)

    # Layer 3
    add_round_rect(slide, (330, 690, 2130, 430), fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.8)
    add_text(slide, "第3层：智能体引擎层（核心）", (395, 720, 360, 38), 18, bold=True)
    add_round_rect(slide, (1080, 695, 600, 74), fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.8)
    add_text(slide, "LangChain + LangGraph（核心大脑）", (1105, 710, 550, 40), 16, bold=True, align=PP_ALIGN.CENTER, wrap=False)

    add_round_rect(slide, (415, 815, 520, 255), fill=WHITE, line=GRID, line_width=1.0)
    left_steps = [
        (430, 835, "需求解析", "方案生成"),
        (430, 915, "成本计算", "专家评审"),
        (430, 995, "辩论协商", "仲裁决策"),
    ]
    for x, y, left_label, right_label in left_steps:
        add_round_rect(slide, (x, y, 155, 54), fill=WHITE, line=GRID, line_width=0.9)
        add_text(slide, left_label, (x + 10, y + 8, 135, 36), 13, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "≫", (598, y + 8, 60, 36), 20, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER, wrap=False)
        add_round_rect(slide, (665, y, 155, 54), fill=WHITE, line=GRID, line_width=0.9)
        add_text(slide, right_label, (675, y + 8, 135, 36), 13, bold=True, align=PP_ALIGN.CENTER)

    engine_cards = [
        (1035, 815, 420, 255, "绿电方案", "核心能力：差分进化算法", "输出：风光储装机配比"),
        (1500, 815, 420, 255, "制冷方案", "核心能力：26种方案库\n温度/水资源修正", "输出：浸没式液冷/间接蒸发等方案"),
        (1965, 815, 420, 255, "供电方案", "核心能力：GB 50174-2017 标准", "输出：A+/A/B/C级架构及配置"),
    ]
    for x, y, w, h, title, body, footer in engine_cards:
        add_round_rect(slide, (x, y, w, h), fill=WHITE, line=GRID, line_width=1.0)
        add_text(slide, title, (x + 30, y + 28, w - 60, 30), 17, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, (x + 38, y + 84, w - 76, 96), 14, color=TEXT_SOFT, align=PP_ALIGN.LEFT)
        add_round_rect(slide, (x + 26, y + 185, w - 52, 48), fill=WHITE, line=GRID, line_width=1.0)
        add_text(slide, footer, (x + 40, y + 192, w - 80, 30), 13, align=PP_ALIGN.CENTER, color=TEXT_SOFT)

    # Layer 4
    add_round_rect(slide, (330, 1148, 2130, 130), fill=WHITE, line=GRID, line_width=1.0)
    add_text(slide, "第4层：工具与算法层", (395, 1178, 320, 34), 18, bold=True)
    tool_items = [
        "差分进化算法",
        "光伏仿真\n(pv_sim)",
        "风电仿真\n(wind_sim)",
        "26种制冷方案库",
        "GB标准供电库",
    ]
    for idx, label in enumerate(tool_items):
        x = 860 + idx * 305
        add_round_rect(slide, (x, 1166, 270, 82), fill=WHITE, line=GRID, line_width=1.0)
        add_text(slide, label, (x + 10, 1182, 250, 48), 14, align=PP_ALIGN.CENTER)

    # Layer 5
    add_round_rect(slide, (330, 1310, 2130, 170), fill=WHITE, line=GRID, line_width=1.0)
    add_text(slide, "第5层：数据层", (395, 1350, 240, 34), 18, bold=True)
    data_items = [
        (900, 1338, 420, 118, "气象数据库\n（温度/光照/风速）"),
        (1600, 1338, 420, 118, "成本数据库\n（设备/运维）"),
        (2200, 1338, 420, 118, "标准数据库\n（规范文档）"),
    ]
    for x, y, w, h, label in data_items:
        add_round_rect(slide, (x, y, w, h), fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.0)
        add_circle(slide, (x + 25, y + 24, 58, 58), fill=GREEN_DARK, line=GREEN_DARK)
        add_text(slide, label, (x + 95, y + 22, w - 120, 74), 14, align=PP_ALIGN.CENTER)


def slide3(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_text(slide, "核心亮点：多智能体协作流", (150, 55, 1120, 88), 27, bold=True)
    add_text(slide, "模拟专家级决策论证过程的 6 阶段闭环", (150, 148, 860, 56), 18, color=TEXT_SOFT)

    add_round_rect(slide, (180, 330, 2460, 300), fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.8)
    top_steps = [
        {
            "asset": assets["p3_icon_1"],
            "asset_box": (110, 350, 200, 200),
            "title_box": (340, 355, 330, 62),
            "desc_box": (355, 442, 260, 100),
            "title": "01. 需求解析",
            "desc": "解析负荷、地点、\n预算、PUE目标",
        },
        {
            "asset": assets["p3_icon_2"],
            "asset_box": (980, 350, 200, 200),
            "title_box": (1210, 355, 350, 62),
            "desc_box": (1225, 442, 285, 112),
            "title": "02. 方案生成",
            "desc": "调度绿电、制冷、供电\n三大底层工具生成初案",
        },
        {
            "asset": assets["p3_icon_3"],
            "asset_box": (1860, 350, 200, 200),
            "title_box": (2090, 355, 300, 62),
            "desc_box": (2105, 442, 255, 102),
            "title": "03. 成本计算",
            "desc": "计算总投资、运营成本与 ROI",
        },
    ]
    for step in top_steps:
        add_picture_contain(slide, step["asset"], step["asset_box"])
        add_round_rect(slide, step["title_box"], fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.6)
        add_text(slide, step["title"], (step["title_box"][0] + 18, step["title_box"][1] + 10, step["title_box"][2] - 36, 40), 18, bold=True, align=PP_ALIGN.CENTER)
        add_round_rect(slide, (step["desc_box"][0] - 16, step["desc_box"][1] - 10, step["desc_box"][2] + 32, step["desc_box"][3] + 20), fill=WHITE, line=GRID, line_width=1.0)
        add_text(slide, step["desc"], step["desc_box"], 15, color=TEXT_SOFT)

    add_text(slide, "→", (760, 382, 70, 42), 24, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER, wrap=False)
    add_text(slide, "→", (1635, 382, 70, 42), 24, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER, wrap=False)

    add_round_rect(slide, (230, 700, 2120, 58), fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.4)
    add_text(slide, "04. 专家评审", (860, 708, 860, 40), 20, bold=True, align=PP_ALIGN.CENTER)

    expert_cards = [
        (420, 790, 700, 112, "经济性分析专家", "评估投资回报率、成本效益"),
        (420, 935, 700, 112, "可靠性分析专家", "评估供电可靠性、冗余配置"),
        (420, 1080, 700, 112, "环保性分析专家", "评估碳排放、绿电消纳率"),
    ]
    for x, y, w, h, title, desc in expert_cards:
        add_round_rect(slide, (x, y, w, h), fill=WHITE, line=GREEN_LINE, line_width=1.2)
        add_circle(slide, (x + 24, y + 24, 56, 56), fill=GREEN_LIGHT)
        add_text(slide, title, (x + 108, y + 18, w - 140, 34), 17, bold=True)
        add_text(slide, desc, (x + 108, y + 52, w - 150, 28), 13, color=TEXT_SOFT)

    add_circle(slide, (1575, 930, 220, 220), fill=GREEN_LIGHT)
    add_text(slide, "↻", (1630, 985, 110, 54), 32, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER, wrap=False)
    add_text(slide, "05. 辩论协商", (1510, 1150, 350, 30), 20, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "多轮论证，观点碰撞，寻找最优平衡点", (1488, 1185, 395, 28), 14, color=TEXT_SOFT, align=PP_ALIGN.CENTER, wrap=False)

    add_picture_contain(slide, assets["p3_gavel"], (2340, 925, 210, 210))
    add_round_rect(slide, (2270, 1080, 300, 66), fill=GREEN_LIGHT, line=GREEN_LINE, line_width=1.4)
    add_text(slide, "06. 仲裁决策", (2295, 1093, 250, 34), 18, bold=True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, (2220, 1162, 375, 95), fill=WHITE, line=GRID, line_width=1.0)
    add_text(slide, "综合各方评分，输出最终\n推荐方案与可行性报告", (2245, 1180, 325, 54), 14, color=TEXT_SOFT, align=PP_ALIGN.CENTER)

    add_round_rect(slide, (520, 1392, 1680, 94), fill=GREEN_LIGHT, line=GRID, line_width=1.0)
    add_text(
        slide,
        "从单一维度计算到多维视角博弈，LangGraph 赋予系统真正的工程级规划智慧。",
        (610, 1412, 1490, 44),
        16,
        color=TEXT_SOFT,
        align=PP_ALIGN.CENTER,
    )


def main() -> None:
    assets = build_assets()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide1(prs, assets)
    slide2(prs)
    slide3(prs, assets)
    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
